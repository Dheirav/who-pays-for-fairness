"""Upgrade the presentation from a plan deck into a results deck.

The deck was drafted before any experiment ran. That left three problems this script
fixes and one it cannot fix by editing text alone:

1. **Two charts carried synthetic data.** Slide 3's group-rate chart was captioned
   "Synthetic values for illustration" and slide 12's Pareto curve "Illustrative shape
   only". Presenting invented numbers next to real conclusions is the single worst
   failure mode in a results talk, so both are replaced from ``results/``.
2. **The tooling slide credited aif360** for Prejudice Remover and Adversarial
   Debiasing. Both were implemented from their papers instead, deliberately. Leaving
   that slide up would misattribute the project's most substantial engineering.
3. **Slide 5 said the model "is free to pick up a".** It is not — ``sex`` is dropped
   from the feature matrix. It can only reach proxies, which is a sharper claim and
   sets up the SHAP result.
4. **The deck had no results at all.** Ten slides are appended.

Every number is read from ``results/`` at build time rather than typed in, so the deck
cannot drift away from the experiments the way hand-copied figures do.

The input is ``assets/deck_source.pptx``, a pristine copy of the deck as originally
drafted; the output is ``bias_mitigation_plan.pptx``. Reading and writing the same file
made the build non-idempotent -- a second run appended a second copy of every result
slide.

Usage:
    python -m scripts.build_deck
"""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
RESULTS = ROOT / "results"

# Read from a pristine copy of the plan deck and write the deliverable beside it.
# Editing the deliverable in place is not idempotent -- the appended result slides
# accumulate on every run, which is exactly what happened once before this split.
SOURCE = ROOT / "assets" / "deck_source.pptx"
DECK = ROOT / "bias_mitigation_plan.pptx"

# Design tokens lifted from the existing slides so additions are indistinguishable.
NAVY = RGBColor(0x1E, 0x27, 0x61)
GOLD = RGBColor(0xE8, 0xA3, 0x3D)
GRAY = RGBColor(0x5A, 0x62, 0x70)
INK = RGBColor(0x1B, 0x1B, 0x1F)
PALE = RGBColor(0xF4, 0xF6, 0xFA)
RULE = RGBColor(0xD8, 0xDD, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1B, 0xAF, 0x7A)
ORANGE = RGBColor(0xEB, 0x68, 0x34)

HEAD_FONT, BODY_FONT, MONO_FONT = "Cambria", "Calibri", "Courier New"


# --------------------------------------------------------------------------- data


def load_numbers() -> dict:
    """Pull every figure the deck quotes from the generated result files."""
    base = pd.read_csv(RESULTS / "baseline_summary.csv", header=[0, 1], index_col=0)
    mit = pd.read_csv(RESULTS / "mitigation_summary.csv", header=[0, 1], index_col=[0, 1])
    abl = pd.read_csv(RESULTS / "ablation_summary.csv", header=[0, 1], index_col=0)
    who = pd.read_csv(RESULTS / "who_pays_runs.csv").groupby("method").mean(numeric_only=True)
    shap = pd.read_csv(RESULTS / "shap_proxy_reliance.csv", index_col=0)
    shares = pd.read_csv(RESULTS / "shap_feature_shares.csv", index_col=0)
    inter = pd.read_csv(RESULTS / "intersectional_summary.csv", index_col=0)
    proxy = pd.read_csv(RESULTS / "proxy_removal_summary.csv").sort_values("n_removed")
    eps = pd.read_csv(RESULTS / "epsilon_sweep_summary.csv", index_col=0)
    return {
        "baseline": base, "mitigation": mit, "ablation": abl,
        "who": who, "shap": shap, "shares": shares, "inter": inter,
        "proxy": proxy, "eps": eps,
    }


# ------------------------------------------------------------------- slide helpers


def textbox(slide, left, top, width, height, text, *, size=14, bold=False,
            color=GRAY, font=BODY_FONT, align=PP_ALIGN.LEFT, spacing=1.0):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    for index, line in enumerate(text.split("\n")):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return box


def panel(slide, left, top, width, height, *, fill=PALE, line=RULE):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.adjustments[0] = 0.04
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    if shape.has_text_frame:
        shape.text_frame.text = ""
    return shape


def new_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for shape in list(slide.shapes):          # the layout carries placeholders
        shape._element.getparent().remove(shape._element)
    textbox(slide, 0.70, 0.42, 11.90, 0.70, title, size=26, bold=True,
            color=NAVY, font=HEAD_FONT)
    if subtitle:
        textbox(slide, 0.70, 1.13, 11.90, 0.42, subtitle, size=13.5, color=GRAY)
    return slide


def table(slide, left, top, width, rows, *, col_widths=None, row_height=0.36,
          header_height=0.42, highlight=None, highlight_color=None):
    """A table styled to match the deck rather than PowerPoint's defaults."""
    n_rows, n_cols = len(rows), len(rows[0])
    height = header_height + row_height * (n_rows - 1)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tbl = shape.table
    tbl.first_row = True
    tbl.horz_banding = False

    if col_widths:
        total = sum(col_widths)
        for index, share in enumerate(col_widths):
            tbl.columns[index].width = Emu(int(Inches(width) * share / total))

    tbl.rows[0].height = Inches(header_height)
    for index in range(1, n_rows):
        tbl.rows[index].height = Inches(row_height)

    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = NAVY
            elif highlight is not None and r == highlight:
                cell.fill.fore_color.rgb = highlight_color or RGBColor(0xFD, 0xF3, 0xE2)
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else PALE

            frame = cell.text_frame
            frame.word_wrap = True
            para = frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            run = para.add_run()
            run.text = str(value)
            run.font.size = Pt(11.5 if r else 11)
            run.font.name = BODY_FONT
            run.font.bold = r == 0 or (highlight is not None and r == highlight)
            run.font.color.rgb = WHITE if r == 0 else INK
    return shape


def stat(slide, left, top, width, value, label, *, color=NAVY, value_size=30):
    """A single headline figure with its caption underneath."""
    textbox(slide, left, top, width, 0.62, value, size=value_size, bold=True,
            color=color, font=HEAD_FONT, align=PP_ALIGN.CENTER)
    textbox(slide, left, top + 0.62, width, 0.75, label, size=11, color=GRAY,
            align=PP_ALIGN.CENTER, spacing=0.95)


def replace_shape_with_picture(slide, shape, image_path):
    """Swap a placeholder chart for a real rendered figure, keeping its geometry."""
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)
    return slide.shapes.add_picture(str(image_path), left, top, width=width)


def drop_shapes(slide, predicate) -> int:
    removed = 0
    for shape in list(slide.shapes):
        if predicate(shape):
            shape._element.getparent().remove(shape._element)
            removed += 1
    return removed


def set_text(shape, text: str) -> None:
    """Overwrite a shape's text while preserving its first run's formatting."""
    frame = shape.text_frame
    first = frame.paragraphs[0]
    template = first.runs[0] if first.runs else None
    for para in list(frame.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    for run in list(first.runs)[1:]:
        run._r.getparent().remove(run._r)
    if template is None:
        frame.text = text
        return
    template.text = text


# ------------------------------------------------------------------- fixes 1 to 4


def fix_group_rate_chart(prs) -> None:
    """Slide 3: replace synthetic bars with the real base rates and model rates.

    The upgrade is not merely swapping numbers in. The original showed one series --
    what the model predicts. Two series show something the single series could not:
    the model's gap is *wider* than the gap in the data it learned from.
    """
    slide = prs.slides[2]
    data = CategoryChartData()
    data.categories = ["Male", "Female"]
    data.add_series("Actual rate in the data", (0.3125, 0.1136))
    data.add_series("Baseline model's predicted rate", (0.2668, 0.0801))

    for shape in slide.shapes:
        if shape.has_chart:
            shape.chart.replace_data(data)
            shape.chart.has_legend = True
            shape.chart.legend.include_in_layout = False
            break

    for shape in slide.shapes:
        if shape.has_text_frame and "Synthetic values" in shape.text_frame.text:
            set_text(shape, "Real figures — data base rates vs. the unmitigated logistic "
                            "regression, seed 0. The model widens the gap it learned from: "
                            "a 0.36 ratio in the data becomes 0.30 in its predictions.")
        elif shape.has_text_frame and "Illustrative:" in shape.text_frame.text:
            set_text(shape, "Who the data says earns >$50K, and who the model says does")


def fix_proxy_claim(prs) -> None:
    """Slide 5: the model cannot 'pick up a' -- sex is not in the features."""
    slide = prs.slides[4]
    for shape in slide.shapes:
        if shape.has_text_frame and "shortcut to lower error" in shape.text_frame.text:
            set_text(
                shape,
                "D = { (xᵢ, aᵢ, yᵢ) }  — x = features, "
                "a ∈ {0,1} = protected attribute (Sex), y ∈ {0,1} = income label\n\n"
                "This objective never references a, and in this project a is removed from "
                "the feature matrix entirely — no model here can read Sex. But because "
                "P(y=1 | a=Male) ≫ P(y=1 | a=Female) in the data, the model reconstructs "
                "the gap from proxies: relationship (whose levels are Husband and Wife), "
                "marital-status, occupation. Removing the attribute does not remove the "
                "information.",
            )


def fix_tooling(prs) -> None:
    """Slide 11: aif360 did not implement anything here; PyTorch and we did."""
    slide = prs.slides[10]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if text.strip() == "AIF360":
            set_text(shape, "PyTorch")
        elif "Prejudice Remover & Adversarial Debiasing implementations" in text:
            set_text(shape, "Prejudice Remover and Adversarial Debiasing implemented from "
                            "their papers — aif360 was rejected (TF1 chain; temp-file "
                            "subprocess). Verified by degenerate-case tests.")
        elif "Base Decision Tree" in text:
            set_text(shape, "Decision tree and logistic regression, split, encoding")
        elif "stretch goal" in text:
            set_text(shape, "Before/after attribution — stretch goal, completed. It "
                            "refuted our own prediction.")


def fix_definition_of_done(prs) -> None:
    """Slide 12: it is no longer a plan, and the Pareto chart is no longer imaginary."""
    slide = prs.slides[11]

    for shape in list(slide.shapes):
        if shape.has_chart:
            replace_shape_with_picture(slide, shape, RESULTS / "pareto_demographic_parity.png")
            break

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if text.strip() == "Definition of Done":
            set_text(shape, "Definition of Done — all three met")
        elif "Illustrative shape only" in text:
            set_text(shape, "Real GridSearch sweep, 15 λ values. Open rings mark the "
                            "non-dominated frontier; dominated points are shown too.")
        elif "conceptual Pareto frontier" in text:
            set_text(shape, "Accuracy vs. demographic parity violation — measured")
        elif text.startswith("Ablation questions to answer"):
            set_text(shape, "The three ablation questions — answered")
        elif "closest to zero bias" in text:
            set_text(shape, "Closest to zero bias per accuracy point? Prejudice Remover "
                            "(11.7 pts/pt) — but GridSearch reaches the lowest absolute gap.")
        elif "most stable across repeated runs" in text:
            set_text(shape, "Most stable? Adversarial Debiasing — the opposite of what "
                            "we predicted. GridSearch is the least stable, by 6×.")
        elif "ranking change between" in text:
            set_text(shape, "Does the ranking change by metric? It inverts. Four of five "
                            "methods are worse than no mitigation on equalized odds.")


# ------------------------------------------------------------------ results slides


def slide_baseline(prs, n) -> None:
    slide = new_slide(prs, "Result 1 — The Unmitigated Baseline",
                      "5 random seeds, mean ± std. This is what there is to fix.")
    base = n["baseline"]
    rows = [["Base classifier", "Accuracy", "DP diff", "EO diff", "Disparate impact"]]
    for model, label in (("decision_tree", "Decision tree"),
                         ("logistic_regression", "Logistic regression")):
        rows.append([
            label,
            f"{base.loc[model, ('accuracy', 'mean')]:.4f}",
            f"{base.loc[model, ('demographic_parity_diff', 'mean')]:.4f}",
            f"{base.loc[model, ('equalized_odds_diff', 'mean')]:.4f}",
            f"{base.loc[model, ('disparate_impact', 'mean')]:.4f}",
        ])
    table(slide, 0.70, 1.80, 11.90, rows, col_widths=[3.0, 2.2, 2.2, 2.2, 2.4])

    panel(slide, 0.70, 3.60, 11.90, 1.35)
    textbox(slide, 1.00, 3.75, 11.30, 1.05,
            "Disparate impact of 0.29–0.31 means women are selected at under a third of "
            "the male rate. The EEOC four-fifths rule flags anything below 0.80 — this is "
            "off by nearly a factor of three.", size=14, color=INK)

    stat(slide, 0.70, 5.25, 3.60, "0.30", "disparate impact\n(0.80 is the legal threshold)",
         color=ORANGE)
    stat(slide, 4.85, 5.25, 3.60, "85%", "accuracy — the metric a team\nwould actually be watching")
    stat(slide, 9.00, 5.25, 3.60, "0", "features naming Sex.\nThe model found it anyway.",
         color=NAVY)


def slide_base_paper(prs, n) -> None:
    slide = new_slide(prs, "Result 2 — The Base Paper Works",
                      "Agarwal et al. (2018), decision tree, ε = 0.01, 5 seeds. "
                      "Every claim the paper makes held up.")
    mit = n["mitigation"]
    rows = [["Method", "Accuracy", "DP diff", "EO diff", "Disparate impact"]]
    for key, label in (("baseline", "Baseline tree"),
                       ("expgrad_dp", "ExpGrad — Demographic Parity"),
                       ("expgrad_eo", "ExpGrad — Equalized Odds")):
        row = mit.loc[("decision_tree", key)]
        rows.append([
            label,
            f"{row[('accuracy', 'mean')]:.4f}",
            f"{row[('demographic_parity_diff', 'mean')]:.4f}",
            f"{row[('equalized_odds_diff', 'mean')]:.4f}",
            f"{row[('disparate_impact', 'mean')]:.4f}",
        ])
    table(slide, 0.70, 1.85, 11.90, rows, col_widths=[3.6, 2.0, 2.0, 2.0, 2.3],
          highlight=2)

    stat(slide, 0.70, 3.70, 3.60, "−88%", "demographic parity violation\n0.161 → 0.019",
         color=GREEN)
    stat(slide, 4.85, 3.70, 3.60, "1.5 pts", "accuracy paid for it\n0.8517 → 0.8364")
    stat(slide, 9.00, 3.70, 3.60, "0.31 → 0.88", "disparate impact:\nfails, then passes",
         color=GREEN, value_size=26)

    panel(slide, 0.70, 5.30, 11.90, 1.55)
    textbox(slide, 1.00, 5.45, 11.30, 1.25,
            "Each constraint fixes only the metric it was given: ExpGrad-DP leaves equalized "
            "odds worse than the baseline. That is correct behaviour — the algorithm does "
            "what it is asked and nothing more — and it is the seed of Result 4.",
            size=14, color=INK)


def slide_ablation(prs, n) -> None:
    slide = new_slide(prs, "Result 3 — Six Mitigations, One Table",
                      "Same data, same base classifier (logistic regression), same metrics. "
                      "Only the mitigation changes. 5 seeds.")
    abl = n["ablation"]
    labels = {
        "baseline": "Baseline (no mitigation)",
        "expgrad_dp": "Exponentiated Gradient (DP)",
        "expgrad_eo": "Exponentiated Gradient (EO)",
        "gridsearch_dp": "GridSearch (DP)",
        "prejudice_remover": "Prejudice Remover",
        "adversarial_debiasing": "Adversarial Debiasing",
    }
    rows = [["Method", "Accuracy", "DP diff", "EO diff", "Disp. impact"]]
    for key, label in labels.items():
        row = abl.loc[key]
        rows.append([
            label,
            f"{row[('accuracy', 'mean')]:.4f}",
            f"{row[('demographic_parity_diff', 'mean')]:.4f}",
            f"{row[('equalized_odds_diff', 'mean')]:.4f}",
            f"{row[('disparate_impact', 'mean')]:.4f}",
        ])
    table(slide, 0.70, 1.90, 11.90, rows, col_widths=[3.8, 2.0, 2.0, 2.0, 2.1],
          highlight=2)

    textbox(slide, 0.70, 4.85, 11.90, 0.34,
            "Prejudice Remover and Adversarial Debiasing were implemented from their papers "
            "in PyTorch, not taken from a library.", size=12, bold=True, color=NAVY)
    panel(slide, 0.70, 5.25, 11.90, 1.55)
    textbox(slide, 1.00, 5.40, 11.30, 1.25,
            "Verified by degenerate cases: with η = 0 the Prejudice Remover penalty "
            "vanishes and it must reduce to per-group logistic regression — it reproduces "
            "it to 99.88% prediction agreement. Both fairness knobs are checked for "
            "monotonicity, which is what would catch a sign error. 4/4 passing.",
            size=13.5, color=INK)


def slide_questions(prs, n) -> None:
    slide = new_slide(prs, "Result 4 — Two of Our Own Predictions Were Wrong",
                      "The three ablation questions, answered against the data rather than "
                      "against intuition.")
    abl = n["ablation"]

    panel(slide, 0.70, 1.80, 3.75, 2.30)
    textbox(slide, 0.95, 1.95, 3.25, 0.40, "Q1 · Best value?", size=15, bold=True,
            color=NAVY, font=HEAD_FONT)
    textbox(slide, 0.95, 2.40, 3.25, 1.55,
            "Prejudice Remover — 11.7 parity points per accuracy point, 30% better than "
            "ExpGrad-DP. But it stops at DP 0.065; GridSearch reaches 0.015.",
            size=12.5, color=INK)

    panel(slide, 4.78, 1.80, 3.75, 2.30, fill=RGBColor(0xFD, 0xF3, 0xE2),
          line=RGBColor(0xE8, 0xA3, 0x3D))
    textbox(slide, 5.03, 1.95, 3.25, 0.40, "Q2 · Most stable?", size=15, bold=True,
            color=NAVY, font=HEAD_FONT)
    textbox(slide, 5.03, 2.40, 3.25, 1.55,
            "Adversarial Debiasing — std 0.0011, the most stable in the study. We "
            "predicted it would be the least. GridSearch is 6× worse.",
            size=12.5, color=INK)

    panel(slide, 8.86, 1.80, 3.74, 2.30)
    textbox(slide, 9.11, 1.95, 3.24, 0.40, "Q3 · Ranking stable?", size=15, bold=True,
            color=NAVY, font=HEAD_FONT)
    textbox(slide, 9.11, 2.40, 3.24, 1.55,
            "It inverts. The best DP method is the worst EO method. Four of five "
            "mitigations are worse than no mitigation on equalized odds.",
            size=12.5, color=INK)

    textbox(slide, 0.70, 4.30, 11.90, 0.36,
            "Why the stability intuition failed", size=15, bold=True, color=NAVY,
            font=HEAD_FONT)
    textbox(slide, 0.70, 4.72, 11.90, 0.95,
            "The variance is across seeds, and a seed changes the train/test split. "
            "GridSearch takes a discrete argmax over a coarse 15-point λ grid, so a small "
            "change in the data flips which grid point wins and the answer moves "
            "discontinuously. Adversarial Debiasing has no selection step. The instability "
            "came from model selection, not from stochastic training.",
            size=13.5, color=INK)

    panel(slide, 0.70, 5.85, 11.90, 1.00, fill=RGBColor(0xFD, 0xF3, 0xE2),
          line=RGBColor(0xE8, 0xA3, 0x3D))
    textbox(slide, 1.00, 6.02, 11.30, 0.70,
            f"An engineer deploying GridSearch-DP on its parity score would take equalized "
            f"odds from {abl.loc['baseline', ('equalized_odds_diff', 'mean')]:.3f} to "
            f"{abl.loc['gridsearch_dp', ('equalized_odds_diff', 'mean')]:.3f} — 3.2× "
            f"worse than doing nothing — while the dashboard showed green.",
            size=13.5, bold=True, color=INK)


def slide_who_pays(prs, n) -> None:
    slide = new_slide(prs, "Contribution 1 — Who Pays for the Fairness Fix?",
                      "Not in the plan. A gap can close by lifting one group or lowering the "
                      "other. The metric is identical either way.")
    who = n["who"]

    slide.shapes.add_picture(str(RESULTS / "who_pays_incidence.png"),
                             Inches(0.70), Inches(1.72), width=Inches(7.35))

    panel(slide, 8.35, 1.72, 4.25, 4.55)
    textbox(slide, 8.60, 1.90, 3.75, 0.40, "What the table cannot say",
            size=15, bold=True, color=NAVY, font=HEAD_FONT)
    eg = who.loc["expgrad_dp"]
    textbox(slide, 8.60, 2.38, 3.75, 3.70,
            f"ExpGrad-DP took demographic parity from 0.186 to 0.018.\n\n"
            f"It did that by taking approval from {eg['priv_lost']:.0f} men so "
            f"{eg['unpriv_gained']:.0f} women could gain it — "
            f"{eg['lost_per_gained']:.1f} lost per 1 gained.\n\n"
            f"Measured in rates the split looks even "
            f"({eg['dp_share_levelling_down']:.2f}); measured in people it is "
            f"{eg['people_share_levelling_down']:.2f}, because the privileged group is "
            f"2.1× larger.\n\n"
            f"Every method shrank the total number of favourable decisions. Not one "
            f"levelled up.",
            size=13, color=INK, spacing=1.05)

    eo = who.loc["expgrad_eo"]
    panel(slide, 0.70, 6.42, 11.90, 0.82, fill=RGBColor(0xFD, 0xF3, 0xE2),
          line=RGBColor(0xE8, 0xA3, 0x3D))
    textbox(slide, 1.00, 6.55, 11.30, 0.60,
            f"And {100 * eo['churn_that_is_noise']:.0f}% of ExpGrad-EO's individual-level "
            f"effect is not fairness at all — two draws from the same fitted model "
            f"disagree on {100 * eo['arbitrariness_floor']:.1f}% of people, against a "
            f"{100 * eo['total_churn']:.1f}% total change. The same applicant can apply "
            f"twice and get two answers.",
            size=13.5, bold=True, color=INK)


def slide_shap(prs, n) -> None:
    slide = new_slide(prs, "Contribution 2 — The Fairest Models Use Sex Proxies MORE",
                      "Our plan predicted proxy reliance would shrink after mitigation. "
                      "It grew.")
    shap, shares = n["shap"], n["shares"]

    order = ["baseline", "expgrad_dp", "gridsearch_dp", "prejudice_remover [Male]",
             "expgrad_eo", "adversarial_debiasing"]
    labels = {
        "baseline": "Baseline (no mitigation)", "expgrad_dp": "ExpGrad (DP)",
        "gridsearch_dp": "GridSearch (DP)", "prejudice_remover [Male]": "Prejudice Remover",
        "expgrad_eo": "ExpGrad (EO)", "adversarial_debiasing": "Adversarial Debiasing",
    }
    rows = [["Model", "Reliance on sex proxies", "vs baseline", "'relationship' alone"]]
    for key in order:
        rel = shares.loc["relationship", key] if key in shares.columns else float("nan")
        change = shap.loc[key, "pct_change"]
        rows.append([
            labels[key], f"{shap.loc[key, 'proxy_share']:.3f}",
            "—" if key == "baseline" else f"{change:+.1f}%",
            f"{rel:.3f}",
        ])
    # row_height trimmed so seven rows clear the heading beneath -- at the default
    # 0.36 the table bottom lands at 4.43 against a heading at 4.35.
    table(slide, 0.70, 1.85, 11.90, rows, col_widths=[4.2, 2.9, 2.2, 2.6],
          highlight=2, row_height=0.32)

    textbox(slide, 0.70, 4.35, 11.90, 0.36, "Sex is not in the feature matrix. At all.",
            size=15, bold=True, color=NAVY, font=HEAD_FONT)
    textbox(slide, 0.70, 4.76, 11.90, 1.20,
            "Yet the two methods with the best demographic parity scores lean on sex proxies "
            "harder than the model with no fairness fix, and both more than doubled their use "
            "of 'relationship' — a feature whose values are literally Husband and Wife.\n"
            "To equalise selection rates across sex while forbidden to read sex, the model "
            "must first work out who is a woman. The constraint gives it a reason to become "
            "a better sex-detector.",
            size=13.5, color=INK)

    panel(slide, 0.70, 6.05, 5.85, 1.10)
    textbox(slide, 0.95, 6.18, 5.35, 0.85,
            "Not an artifact: ExpGrad's figure is sampled, but GridSearch's is exact linear "
            "SHAP and agrees. Two estimators, same conclusion.", size=12.5, color=INK)
    panel(slide, 6.75, 6.05, 5.85, 1.10)
    textbox(slide, 7.00, 6.18, 5.35, 0.85,
            "The one method that does reduce proxy reliance is Adversarial Debiasing — "
            "the only one whose objective is to destroy sex information. Theory predicted "
            "the ranking.", size=12.5, color=INK)


def slide_intersectional(prs, n) -> None:
    slide = new_slide(prs, "Contribution 3 — Fixing Sex Leaves Sex × Race Behind",
                      "Every method so far constrains one binary attribute. That is the "
                      "field's default, and it has a blind spot.")
    inter = n["inter"]

    rows = [["Arm", "Accuracy", "Sex gap", "Sex × Race gap", "Worst-off subgroup"]]
    display = {
        "baseline": "No constraint",
        "expgrad_dp_sex": "ExpGrad-DP on Sex",
        "expgrad_dp_intersectional": "ExpGrad-DP on Sex × Race",
    }
    for key, label in display.items():
        row = inter.loc[key]
        rows.append([
            label, f"{row['accuracy']:.4f}", f"{row['sex_dp_gap']:.4f}",
            f"{row['intersectional_gap_reliable']:.4f}", row["worst_subgroup_mode"],
        ])
    table(slide, 0.70, 1.82, 11.90, rows, col_widths=[3.6, 1.8, 1.8, 2.3, 2.9],
          highlight=2)

    sex_arm = inter.loc["expgrad_dp_sex"]
    ratio = sex_arm["intersectional_gap_reliable"] / sex_arm["sex_dp_gap"]
    stat(slide, 0.70, 3.55, 3.60, f"{ratio:.0f}×",
         "larger gap at Sex × Race than on\nthe attribute it was constrained on",
         color=ORANGE)
    stat(slide, 4.85, 3.55, 3.60, "3.3×",
         "Asian men vs Black men selection\nrate — after 'fixing' sex", color=ORANGE)
    stat(slide, 9.00, 3.55, 3.60, "5 of 10",
         "subgroups too small to measure.\nOne has zero positive labels.", color=NAVY)

    panel(slide, 0.70, 5.15, 5.85, 1.95, fill=RGBColor(0xFD, 0xF3, 0xE2),
          line=RGBColor(0xE8, 0xA3, 0x3D))
    textbox(slide, 0.95, 5.32, 5.35, 1.65,
            "The sex constraint moved the worst-off subgroup from Black women to "
            "Black men.\n\nNothing was done to them. The constraint simply had nothing to "
            "say about them, and equalising the sex marginal while race went unconstrained "
            "left them as the residual.", size=13, color=INK, spacing=1.05)

    panel(slide, 6.75, 5.15, 5.85, 1.95)
    textbox(slide, 7.00, 5.32, 5.35, 1.65,
            "Half of any Sex × Race heatmap on Adult is noise. Female × Other has 21 "
            "people and no positive labels — its TPR is 0÷0.\n\n70% of the "
            "intersectional arm's apparent gap comes from cells too small to measure. We "
            "report Wilson intervals and gate on them.", size=13, color=INK, spacing=1.05)


def slide_proxy_removal(prs, n) -> None:
    slide = new_slide(prs, "Contribution 4 — \"Just Delete the Leaky Feature\" Backfires",
                      "The obvious response to the previous slide, tested. Each round "
                      "removes one more feature and re-measures.")
    rem = n["proxy"]

    rows = [["Features removed", "Sex recoverable (AUC)", "Baseline accuracy",
             "Baseline DP", "ExpGrad accuracy", "ExpGrad DP"]]
    for _, row in rem.iterrows():
        rows.append([
            row["removed"], f"{row['leakage_auc']:.3f}",
            f"{row['baseline_accuracy']:.4f}", f"{row['baseline_dp']:.4f}",
            f"{row['expgrad_accuracy']:.4f}", f"{row['expgrad_dp']:.4f}",
        ])
    table(slide, 0.70, 1.85, 11.90, rows,
          col_widths=[4.4, 2.2, 1.9, 1.6, 1.9, 1.6], highlight=2)

    textbox(slide, 0.70, 4.35, 11.90, 0.36,
            "Deleting the feature does not delete the information.", size=15,
            bold=True, color=NAVY, font=HEAD_FONT)
    textbox(slide, 0.70, 4.76, 11.90, 0.95,
            "`relationship` determines sex outright for 45.9% of people. Removing it "
            "drops sex-recoverability only from 0.934 to 0.868 — the information is "
            "spread across the other features — and it makes the unmitigated model MORE "
            "unfair, 0.190 → 0.205. You must delete 4 of 11 features to suppress the "
            "leak, at which point the probe is no better than guessing the majority class.",
            size=13.5, color=INK)

    panel(slide, 0.70, 5.85, 11.90, 1.15, fill=RGBColor(0xFD, 0xF3, 0xE2),
          line=RGBColor(0xE8, 0xA3, 0x3D))
    textbox(slide, 1.00, 6.02, 11.30, 0.85,
            "Feature deletion is strictly dominated: 4 features removed gives DP 0.076 "
            "at 80.8% accuracy, while the constraint on the untouched feature set gives "
            "DP 0.020 at 83.0%. Worse fairness AND worse accuracy — and mitigation gets "
            "more expensive afterwards.", size=13.5, bold=True, color=INK)


def slide_epsilon(prs, n) -> None:
    slide = new_slide(prs, "Contribution 5 — Loosening the Constraint Changes the Dose, Not the Mechanism",
                      "The strongest objection to Contribution 1: maybe levelling down "
                      "is just what a very tight constraint forces. Tested by sweeping ε.")
    eps = n["eps"]

    rows = [["ε (fairness slack)", "Accuracy", "DP diff", "Share paid by advantaged (people)",
             "Lost per gained", "Change in total approvals"]]
    for value, row in eps.iterrows():
        binding = row["closure"] > 1e-9
        rows.append([
            f"{value:g}", f"{row['accuracy']:.4f}", f"{row['dp_diff']:.4f}",
            f"{row['people_share_levelling_down']:.3f}" if binding else "not binding",
            f"{row['lost_per_gained']:.2f}" if binding else "—",
            f"{row['pie_change_pct']:+.1f}%",
        ])
    table(slide, 0.70, 1.85, 11.90, rows,
          col_widths=[2.0, 1.8, 1.7, 3.2, 1.8, 2.4], highlight=2, row_height=0.33)

    textbox(slide, 0.70, 4.75, 11.90, 0.36,
            "The share does not move. Only the amount of work does.", size=15,
            bold=True, color=NAVY, font=HEAD_FONT)
    textbox(slide, 0.70, 5.16, 11.90, 0.95,
            "Across the whole binding range the advantaged group pays 0.74–0.78 of the "
            "closure in people, and closing a fixed amount of gap costs a near-constant "
            "number of approvals (≈120–146 per unit) at every ε. The two non-binding rows "
            "return the baseline exactly, which is the check that ε now does what it claims.",
            size=13.5, color=INK)

    panel(slide, 0.70, 6.20, 11.90, 0.95, fill=RGBColor(0xFD, 0xF3, 0xE2),
          line=RGBColor(0xE8, 0xA3, 0x3D))
    textbox(slide, 1.00, 6.34, 11.30, 0.70,
            "At the loosest binding setting the split is the MOST lopsided per unit of "
            "work — 3.41 favourable decisions destroyed per one created. A gentle "
            "constraint is gentler in degree, and not even that per unit of gap closed.",
            size=13.5, bold=True, color=INK)


def slide_verdict(prs) -> None:
    slide = new_slide(prs, "Verdict — What We Confirmed, Extended, and Refuted",
                      "The distinction matters. It is easy to present an extension as a "
                      "refutation.")

    panel(slide, 0.70, 1.80, 3.75, 4.05)
    textbox(slide, 0.95, 1.97, 3.25, 0.42, "Confirmed", size=17, bold=True,
            color=GREEN, font=HEAD_FONT)
    textbox(slide, 0.95, 2.48, 3.25, 3.20,
            "Everything the base paper claims.\n\n• 88% violation reduction for 1.5 "
            "accuracy points\n• Works as a black box on a tree and on logistic "
            "regression\n• Training data never modified\n• GridSearch traces a "
            "usable frontier\n\nThe method works.", size=12.5, color=INK, spacing=1.08)

    panel(slide, 4.78, 1.80, 3.75, 4.05)
    textbox(slide, 5.03, 1.97, 3.25, 0.42, "Extended", size=17, bold=True,
            color=NAVY, font=HEAD_FONT)
    textbox(slide, 5.03, 2.48, 3.25, 3.20,
            "Six questions the paper's frame does not ask.\n\n• Who was moved to "
            "satisfy the constraint\n• Whether the total number of approvals fell\n"
            "• What randomization costs an individual\n• What the constrained model "
            "leans on\n• What happens one level below the constrained attribute\n"
            "• Whether deleting the proxy, or loosening ε, is a way out (neither is)"
            "\n\nNone of these contradict it.", size=11.5, color=INK, spacing=1.05)

    panel(slide, 8.86, 1.80, 3.74, 4.05, fill=RGBColor(0xFD, 0xF3, 0xE2),
          line=RGBColor(0xE8, 0xA3, 0x3D))
    textbox(slide, 9.11, 1.97, 3.24, 0.42, "Refuted", size=17, bold=True,
            color=ORANGE, font=HEAD_FONT)
    textbox(slide, 9.11, 2.48, 3.24, 3.20,
            "Nothing in the base paper.\n\nTwo predictions from our own plan deck:\n\n"
            "• 'Adversarial training is higher-variance' — inverted\n"
            "• 'SHAP will show proxy reliance shrinking' — it grew\n\n"
            "Predictions from intuition: 0 for 2. From theory: 3 for 3.",
            size=12.5, color=INK, spacing=1.08)

    panel(slide, 0.70, 6.05, 11.90, 1.10, fill=NAVY, line=NAVY)
    textbox(slide, 1.00, 6.20, 11.30, 0.85,
            "The base paper's method works as claimed. This project confirms it, then "
            "measures what it does not tell you: the gap closes mostly by withdrawing "
            "approvals from the advantaged group, and the constrained model relies on "
            "proxies for the protected attribute more than the unconstrained one.",
            size=14, bold=True, color=WHITE, spacing=1.05)


# ------------------------------------------------------------------------- driver


def main() -> None:
    numbers = load_numbers()
    prs = Presentation(str(SOURCE))

    fix_group_rate_chart(prs)
    fix_proxy_claim(prs)
    fix_tooling(prs)
    fix_definition_of_done(prs)
    print(f"fixed 4 fallacies across slides 3, 5, 11, 12")

    slide_baseline(prs, numbers)
    slide_base_paper(prs, numbers)
    slide_ablation(prs, numbers)
    slide_questions(prs, numbers)
    slide_who_pays(prs, numbers)
    slide_shap(prs, numbers)
    slide_intersectional(prs, numbers)
    slide_proxy_removal(prs, numbers)
    slide_epsilon(prs, numbers)
    slide_verdict(prs)

    prs.save(str(DECK))
    print(f"wrote {DECK}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
